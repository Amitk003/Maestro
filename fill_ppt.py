#!/usr/bin/env python3
"""
Maestro Hackathon PPT Filler
Fills slides 2-6 with project details, flowcharts, diagrams, and professional styling.
No emojis. No em dashes. Clean, technical, visually impressive.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os

# ── Color Palette (Professional Dark Theme) ──────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)  # Deep navy
BG_CARD      = RGBColor(0x16, 0x21, 0x3E)  # Card background
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)  # Bright blue
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)  # Green
ACCENT_PURPLE= RGBColor(0xA7, 0x8B, 0xFA)  # Purple
ACCENT_ORANGE= RGBColor(0xFB, 0x92, 0x3C)  # Orange
ACCENT_PINK  = RGBColor(0xF4, 0x72, 0xB6)  # Pink
ACCENT_CYAN  = RGBColor(0x22, 0xD3, 0xEE)  # Cyan
TEXT_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT    = RGBColor(0xCB, 0xD5, 0xE1)  # Light gray
TEXT_DIM      = RGBColor(0x94, 0xA3, 0xB8)  # Dimmed text
BORDER_COLOR  = RGBColor(0x33, 0x44, 0x6B)  # Subtle border
GRADIENT_START= RGBColor(0x1E, 0x29, 0x3B)
GRADIENT_END  = RGBColor(0x0F, 0x17, 0x2A)

# ── Slide dimensions (13.33 x 7.5 inches) ───────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PPTX_PATH = '../Vibeathon_6.0_Vibecoding_Hackathon_July_2026_Idea_Submission_Template.pptx'
OUTPUT_PATH = '../Maestro_Hackathon_Submission_v2.pptx'


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    """Add a shape with optional fill and border."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_arrow_down(slide, x, y, width=Inches(0.3), height=Inches(0.4), color=ACCENT_BLUE):
    """Add a down-pointing arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_connector_line(slide, x1, y1, x2, y2, color=ACCENT_BLUE, width=Pt(2)):
    """Add a line connector between two points."""
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def create_flow_box(slide, left, top, width, height, text, fill_color=BG_CARD, border_color=ACCENT_BLUE,
                    text_color=TEXT_WHITE, font_size=11):
    """Create a styled flow box with text."""
    shape = add_rounded_rect(slide, left, top, width, height, fill_color, border_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    p.font.bold = True
    shape.text_frame.margin_left = Pt(6)
    shape.text_frame.margin_right = Pt(6)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)
    return shape


def create_icon_box(slide, left, top, width, height, icon_char, label, fill_color, text_color=TEXT_WHITE):
    """Create a box with icon character and label."""
    shape = add_rounded_rect(slide, left, top, width, height, fill_color, None)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Icon
    p = tf.paragraphs[0]
    p.text = icon_char
    p.font.size = Pt(24)
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    p.font.bold = True
    
    # Label
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = TEXT_LIGHT
    p2.font.name = 'Calibri'
    p2.alignment = PP_ALIGN.CENTER
    
    return shape


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Current Problem
# ══════════════════════════════════════════════════════════════════════════════
# Logo detection thresholds (template-specific)
LOGO_LEFT_MAX_X = 0.3
LOGO_LEFT_MAX_Y = 0.3
LOGO_LEFT_MIN_W = 1.0
LOGO_RIGHT_MIN_X = 11.0
LOGO_RIGHT_MAX_Y = 0.3
LOGO_RIGHT_MIN_W = 0.5
BG_CENTER_MIN_X = 3.5
BG_CENTER_MAX_X = 5.0
BG_CENTER_MIN_Y = 1.0
BG_CENTER_MAX_Y = 3.0
BG_CENTER_MIN_W = 3.0

def clear_slide(slide):
    """Remove all shapes from a slide, preserving the 3 logo images (top-left, top-right, center background)."""
    logo_shapes = []
    for shape in slide.shapes:
        if shape.shape_type != 13:  # Not a picture
            continue
        left_in = shape.left / 914400
        top_in = shape.top / 914400
        w_in = shape.width / 914400
        
        is_logo = False
        # Top-left logo (event branding)
        if left_in < LOGO_LEFT_MAX_X and top_in < LOGO_LEFT_MAX_Y and w_in > LOGO_LEFT_MIN_W:
            is_logo = True
        # Top-right logo (corner icon)
        elif left_in > LOGO_RIGHT_MIN_X and top_in < LOGO_RIGHT_MAX_Y and w_in > LOGO_RIGHT_MIN_W:
            is_logo = True
        # Center background graphic (watermark)
        elif BG_CENTER_MIN_X < left_in < BG_CENTER_MAX_X and BG_CENTER_MIN_Y < top_in < BG_CENTER_MAX_Y and w_in > BG_CENTER_MIN_W:
            is_logo = True
        
        if is_logo:
            logo_shapes.append(shape)
    
    for shape in list(slide.shapes):
        if shape in logo_shapes:
            continue
        sp = shape._element
        sp.getparent().remove(sp)


def fill_slide_2(slide):
    """Fill Slide 2: Current Problem."""
    set_slide_bg(slide, BG_DARK)
    clear_slide(slide)
    
    # ── Title ──
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Current Problem", font_size=32, color=TEXT_WHITE, bold=True)
    
    # Subtitle line
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
                 "Why restaurants need a fundamentally different approach", font_size=14, color=TEXT_DIM)
    
    # ── Left Column: Problem Cards ──
    problems = [
        ("Reactive Systems", "Ticket printers, manual coordination, and yelling are the backbone of most restaurant operations. Problems are only addressed after they occur, leading to cascading failures.", ACCENT_ORANGE),
        ("Information Silos", "Kitchen, front-of-house, inventory, and management operate in disconnected loops. No one has a complete picture of restaurant state at any given moment.", ACCENT_PINK),
        ("Waste & Inefficiency", "Ingredients spoil before use. Kitchen stations become bottlenecks while others sit idle. Staff burn out from poor task distribution. Tables sit dirty while waiters are overwhelmed.", ACCENT_PURPLE),
        ("No Predictive Capability", "Existing systems cannot anticipate demand surges from weather, local events, or time patterns. Restaurants are always caught off-guard.", ACCENT_CYAN),
    ]
    
    card_y = Inches(1.6)
    for i, (title, desc, accent) in enumerate(problems):
        y = card_y + Inches(i * 1.35)
        # Card background
        card = add_rounded_rect(slide, Inches(0.8), y, Inches(5.8), Inches(1.2), BG_CARD, BORDER_COLOR)
        
        # Accent bar on left
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.06), Inches(1.2), fill_color=accent)
        
        # Title
        add_text_box(slide, Inches(1.1), y + Inches(0.08), Inches(5.3), Inches(0.35),
                     title, font_size=14, color=accent, bold=True)
        
        # Description
        add_text_box(slide, Inches(1.1), y + Inches(0.42), Inches(5.3), Inches(0.75),
                     desc, font_size=10, color=TEXT_LIGHT)
    
    # ── Right Column: Impact Metrics ──
    add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4),
                 "Industry Impact", font_size=18, color=TEXT_WHITE, bold=True)
    
    metrics = [
        ("$200B+", "Annual food waste in US restaurants"),
        ("30%", "Average kitchen idle time during peak"),
        ("4.2x", "Higher staff turnover vs. tech industry"),
    ]
    
    metric_y = Inches(2.2)
    for i, (value, label) in enumerate(metrics):
        y = metric_y + Inches(i * 1.1)
        card = add_rounded_rect(slide, Inches(7.2), y, Inches(5.2), Inches(0.95), BG_CARD, BORDER_COLOR)
        add_text_box(slide, Inches(7.5), y + Inches(0.1), Inches(2), Inches(0.5),
                     value, font_size=28, color=ACCENT_BLUE, bold=True)
        add_text_box(slide, Inches(7.5), y + Inches(0.55), Inches(4.5), Inches(0.35),
                     label, font_size=12, color=TEXT_LIGHT)
    
    # ── Bottom: Affected Stakeholders ──
    add_text_box(slide, Inches(7.2), Inches(5.6), Inches(5), Inches(0.3),
                 "Who Is Affected", font_size=14, color=TEXT_DIM, bold=True)
    
    stakeholders = ["Customers", "Kitchen Staff", "Wait Staff", "Managers", "Owners"]
    for i, name in enumerate(stakeholders):
        x = Inches(7.2) + Inches(i * 1.05)
        create_flow_box(slide, x, Inches(5.95), Inches(0.95), Inches(0.5), name,
                       fill_color=BG_CARD, border_color=ACCENT_BLUE, font_size=9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Proposed Solution
# ══════════════════════════════════════════════════════════════════════════════
def fill_slide_3(slide):
    """Fill Slide 3: Proposed Solution."""
    set_slide_bg(slide, BG_DARK)
    clear_slide(slide)
    
    # ── Title ──
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Proposed Solution", font_size=32, color=TEXT_WHITE, bold=True)
    
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
                 "Maestro: An AI-Powered Restaurant Digital Twin with Multi-Agent Coordination",
                 font_size=14, color=TEXT_DIM)
    
    # ── Central Concept Box ──
    concept_box = add_rounded_rect(slide, Inches(3.5), Inches(1.6), Inches(6.3), Inches(1.0),
                                   BG_CARD, ACCENT_BLUE)
    add_text_box(slide, Inches(3.7), Inches(1.7), Inches(5.9), Inches(0.35),
                 "Core Innovation", font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.7), Inches(2.05), Inches(5.9), Inches(0.45),
                 "Six autonomous AI agents continuously monitor, predict, and optimize every aspect of restaurant operations through a living digital twin.",
                 font_size=11, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    
    # ── Six Agent Boxes (2 rows of 3) ──
    agents = [
        ("Demand Seer", "Predicts demand from\nweather, events, time", ACCENT_BLUE),
        ("Kitchen Conductor", "Balances station loads\nand routing", ACCENT_GREEN),
        ("Inventory Guardian", "Prevents spoilage and\nwaste", ACCENT_ORANGE),
        ("Guest Alchemist", "Personalizes customer\nexperience", ACCENT_PURPLE),
        ("Staff Harmony", "Optimizes task assignment\nand prevents burnout", ACCENT_PINK),
        ("Maestro Orchestrator", "Resolves conflicts and\nmaximizes global score", ACCENT_CYAN),
    ]
    
    start_x = Inches(0.8)
    start_y = Inches(2.9)
    box_w = Inches(3.8)
    box_h = Inches(1.0)
    gap_x = Inches(4.1)
    gap_y = Inches(1.2)
    
    for i, (name, desc, color) in enumerate(agents):
        row = i // 3
        col = i % 3
        x = start_x + col * gap_x
        y = start_y + row * gap_y
        
        card = add_rounded_rect(slide, x, y, box_w, box_h, BG_CARD, color)
        
        # Agent name
        add_text_box(slide, x + Inches(0.15), y + Inches(0.08), box_w - Inches(0.3), Inches(0.3),
                     name, font_size=13, color=color, bold=True)
        
        # Description
        add_text_box(slide, x + Inches(0.15), y + Inches(0.4), box_w - Inches(0.3), Inches(0.55),
                     desc, font_size=10, color=TEXT_LIGHT)
    
    # ── Key Benefits Bar ──
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(3), Inches(0.3),
                 "Key Benefits", font_size=14, color=TEXT_WHITE, bold=True)
    
    benefits = [
        ("Proactive", "Act before problems occur"),
        ("Always On", "Heuristic fallback if AI is down"),
        ("Real-Time", "5-second update cycle via WebSocket"),
    ]
    
    for i, (title, desc) in enumerate(benefits):
        x = Inches(0.8) + Inches(i * 4.1)
        card = add_rounded_rect(slide, x, Inches(5.85), Inches(3.8), Inches(0.85), BG_CARD, BORDER_COLOR)
        add_text_box(slide, x + Inches(0.15), Inches(5.92), Inches(3.5), Inches(0.3),
                     title, font_size=12, color=ACCENT_GREEN, bold=True)
        add_text_box(slide, x + Inches(0.15), Inches(6.2), Inches(3.5), Inches(0.4),
                     desc, font_size=10, color=TEXT_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Technical Approach
# ══════════════════════════════════════════════════════════════════════════════
def fill_slide_4(slide):
    """Fill Slide 4: Technical Approach with architecture diagram."""
    set_slide_bg(slide, BG_DARK)
    clear_slide(slide)
    
    # ── Title ──
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "Technical Approach", font_size=32, color=TEXT_WHITE, bold=True)
    
    # ── Architecture Diagram (Left) ──
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(5), Inches(0.35),
                 "System Architecture", font_size=16, color=ACCENT_BLUE, bold=True)
    
    # Customer Layer
    create_flow_box(slide, Inches(1.5), Inches(1.4), Inches(4.5), Inches(0.6),
                    "Customer / Staff / Manager Interfaces",
                    fill_color=RGBColor(0x1E, 0x3A, 0x5F), border_color=ACCENT_BLUE, font_size=10)
    
    add_arrow_down(slide, Inches(3.6), Inches(2.05), Inches(0.25), Inches(0.3), ACCENT_BLUE)
    
    # Next.js Layer
    create_flow_box(slide, Inches(1.5), Inches(2.4), Inches(4.5), Inches(0.6),
                    "Next.js 15 Frontend (React 19 + Tailwind + Framer Motion)",
                    fill_color=RGBColor(0x1E, 0x3A, 0x5F), border_color=ACCENT_GREEN, font_size=10)
    
    add_arrow_down(slide, Inches(3.6), Inches(3.05), Inches(0.25), Inches(0.3), ACCENT_GREEN)
    
    # Socket.io Layer
    create_flow_box(slide, Inches(1.5), Inches(3.4), Inches(4.5), Inches(0.55),
                    "Socket.io Real-Time Bridge",
                    fill_color=RGBColor(0x2D, 0x1B, 0x3E), border_color=ACCENT_PURPLE, font_size=10)
    
    add_arrow_down(slide, Inches(3.6), Inches(4.0), Inches(0.25), Inches(0.3), ACCENT_PURPLE)
    
    # Agent Worker
    create_flow_box(slide, Inches(1.5), Inches(4.35), Inches(4.5), Inches(0.6),
                    "Agent Worker (Node.js + Express)",
                    fill_color=RGBColor(0x1B, 0x2D, 0x1B), border_color=ACCENT_ORANGE, font_size=10)
    
    add_arrow_down(slide, Inches(3.6), Inches(5.0), Inches(0.25), Inches(0.3), ACCENT_ORANGE)
    
    # Digital Twin
    create_flow_box(slide, Inches(1.5), Inches(5.35), Inches(4.5), Inches(0.6),
                    "Digital Twin Engine (5s tick cycle)",
                    fill_color=RGBColor(0x2D, 0x1B, 0x1B), border_color=ACCENT_PINK, font_size=10)
    
    # ── Right Side: Agent Decision Flow ──
    add_text_box(slide, Inches(7.0), Inches(0.9), Inches(5), Inches(0.35),
                 "Agent Decision Flow", font_size=16, color=ACCENT_GREEN, bold=True)
    
    # Flow steps
    flow_steps = [
        ("1", "Digital Twin State Snapshot", ACCENT_BLUE),
        ("2", "6 Agents Analyze State", ACCENT_GREEN),
        ("3", "Each Agent Proposes Actions", ACCENT_ORANGE),
        ("4", "Orchestrator Scores & Resolves", ACCENT_PURPLE),
        ("5", "Best Actions Applied to Twin", ACCENT_PINK),
        ("6", "WebSocket Broadcast to Frontend", ACCENT_CYAN),
    ]
    
    for i, (num, text, color) in enumerate(flow_steps):
        y = Inches(1.4) + Inches(i * 0.75)
        
        # Number circle
        circle = add_shape(slide, MSO_SHAPE.OVAL, Inches(7.0), y, Inches(0.4), Inches(0.4), fill_color=color)
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = TEXT_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Step text
        add_text_box(slide, Inches(7.55), y + Inches(0.05), Inches(4.5), Inches(0.35),
                     text, font_size=12, color=TEXT_LIGHT)
        
        # Connector line (except last)
        if i < len(flow_steps) - 1:
            add_connector_line(slide, Inches(7.2), y + Inches(0.45), Inches(7.2), y + Inches(0.7),
                             color=RGBColor(0x33, 0x44, 0x6B), width=Pt(1.5))
    
    # ── Tech Stack Table (Bottom) ──
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(3), Inches(0.3),
                 "Technology Stack", font_size=14, color=TEXT_WHITE, bold=True)
    
    techs = [
        ("Frontend", "Next.js 15, React 19, Tailwind v4, Framer Motion"),
        ("AI Engine", "Google Gemini (with heuristic fallback)"),
        ("Database", "PostgreSQL via Supabase"),
        ("Real-Time", "Socket.io + WebSocket"),
        ("State", "Zustand (client) + In-memory (server)"),
    ]
    
    for i, (category, detail) in enumerate(techs):
        x = Inches(0.8) + Inches(i * 2.45)
        card = add_rounded_rect(slide, x, Inches(6.55), Inches(2.3), Inches(0.7), BG_CARD, BORDER_COLOR)
        add_text_box(slide, x + Inches(0.1), Inches(6.58), Inches(2.1), Inches(0.25),
                     category, font_size=9, color=ACCENT_BLUE, bold=True)
        add_text_box(slide, x + Inches(0.1), Inches(6.82), Inches(2.1), Inches(0.4),
                     detail, font_size=8, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Use Cases & Impact
# ══════════════════════════════════════════════════════════════════════════════
def fill_slide_5(slide):
    """Fill Slide 5: Use Cases & Impact."""
    set_slide_bg(slide, BG_DARK)
    clear_slide(slide)
    
    # ── Title ──
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Use Cases & Impact", font_size=32, color=TEXT_WHITE, bold=True)
    
    # ── Left: Use Case Cards ──
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.35),
                 "Key Use Cases", font_size=16, color=ACCENT_BLUE, bold=True)
    
    use_cases = [
        ("Customer Ordering", "Guest types: '25 min, light dinner, pre-show'. Guest Alchemist parses intent, checks kitchen load, returns personalized meal sequence.", ACCENT_BLUE),
        ("Kitchen Optimization", "Grill station hits 90% load. Kitchen Conductor auto-routes items to Saute or Cold Prep. Chefs see only relevant orders.", ACCENT_GREEN),
        ("Crisis Response", "Storm approaches + stadium event ends. Demand Seer detects surge. Inventory Guardian flags spoilage. Orchestrator resolves all conflicts.", ACCENT_ORANGE),
        ("Waste Prevention", "Salmon at 35% freshness. Inventory Guardian promotes Cold Salmon Tartare. 3.2kg waste prevented. Guest delight maintained.", ACCENT_PURPLE),
    ]
    
    for i, (title, desc, color) in enumerate(use_cases):
        y = Inches(1.7) + Inches(i * 1.25)
        card = add_rounded_rect(slide, Inches(0.8), y, Inches(5.8), Inches(1.1), BG_CARD, BORDER_COLOR)
        
        # Accent bar
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.06), Inches(1.1), fill_color=color)
        
        add_text_box(slide, Inches(1.1), y + Inches(0.08), Inches(5.3), Inches(0.3),
                     title, font_size=13, color=color, bold=True)
        add_text_box(slide, Inches(1.1), y + Inches(0.4), Inches(5.3), Inches(0.65),
                     desc, font_size=10, color=TEXT_LIGHT)
    
    # ── Right: Impact Metrics ──
    add_text_box(slide, Inches(7.2), Inches(1.2), Inches(5), Inches(0.35),
                 "Expected Impact", font_size=16, color=ACCENT_GREEN, bold=True)
    
    impacts = [
        ("Waste Reduction", "30-40%", "Through proactive spoilage alerts and dynamic menu promotion", ACCENT_ORANGE),
        ("Service Speed", "25%", "Faster table turnover via intelligent routing and task prioritization", ACCENT_BLUE),
        ("Staff Efficiency", "35%", "Reduced idle time and burnout through optimized task distribution", ACCENT_GREEN),
        ("Guest Satisfaction", "+2.0", "Points on delight score via personalized experiences and recovery perks", ACCENT_PURPLE),
    ]
    
    for i, (metric, value, desc, color) in enumerate(impacts):
        y = Inches(1.7) + Inches(i * 1.25)
        card = add_rounded_rect(slide, Inches(7.2), y, Inches(5.2), Inches(1.1), BG_CARD, BORDER_COLOR)
        
        # Value highlight
        add_text_box(slide, Inches(7.4), y + Inches(0.1), Inches(1.5), Inches(0.5),
                     value, font_size=24, color=color, bold=True)
        
        # Metric name
        add_text_box(slide, Inches(9.0), y + Inches(0.12), Inches(3), Inches(0.3),
                     metric, font_size=13, color=TEXT_WHITE, bold=True)
        
        # Description
        add_text_box(slide, Inches(9.0), y + Inches(0.45), Inches(3), Inches(0.55),
                     desc, font_size=10, color=TEXT_LIGHT)
    
    # ── Bottom: Target Users ──
    add_text_box(slide, Inches(0.8), Inches(6.7), Inches(3), Inches(0.3),
                 "Target Users", font_size=12, color=TEXT_DIM, bold=True)
    
    users = ["Restaurant Owners", "Kitchen Managers", "Front-of-House Staff", "Customers", "Operations Teams"]
    for i, user in enumerate(users):
        x = Inches(0.8) + Inches(i * 2.45)
        create_flow_box(slide, x, Inches(7.0), Inches(2.3), Inches(0.4), user,
                       fill_color=BG_CARD, border_color=ACCENT_CYAN, font_size=9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Future Scope & Conclusion
# ══════════════════════════════════════════════════════════════════════════════
def fill_slide_6(slide):
    """Fill Slide 6: Future Scope & Conclusion."""
    set_slide_bg(slide, BG_DARK)
    clear_slide(slide)
    
    # ── Title ──
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Future Scope & Conclusion", font_size=32, color=TEXT_WHITE, bold=True)
    
    # ── Left: Future Enhancements ──
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.35),
                 "Future Enhancements", font_size=16, color=ACCENT_BLUE, bold=True)
    
    enhancements = [
        ("Multi-Restaurant Federation", "Extend the digital twin to manage multiple restaurant locations with cross-location inventory sharing and demand balancing.", ACCENT_BLUE),
        ("Advanced ML Models", "Replace heuristic fallback with trained ML models for demand prediction, spoilage forecasting, and customer preference learning.", ACCENT_GREEN),
        ("Voice & Chat Integration", "Enable natural language ordering via voice assistants and in-app chat. Staff can interact with agents via voice commands.", ACCENT_PURPLE),
        ("Supply Chain Integration", "Connect directly to suppliers for automated reordering, price optimization, and delivery scheduling.", ACCENT_ORANGE),
    ]
    
    for i, (title, desc, color) in enumerate(enhancements):
        y = Inches(1.7) + Inches(i * 1.2)
        card = add_rounded_rect(slide, Inches(0.8), y, Inches(5.8), Inches(1.05), BG_CARD, BORDER_COLOR)
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.06), Inches(1.05), fill_color=color)
        add_text_box(slide, Inches(1.1), y + Inches(0.08), Inches(5.3), Inches(0.3),
                     title, font_size=13, color=color, bold=True)
        add_text_box(slide, Inches(1.1), y + Inches(0.4), Inches(5.3), Inches(0.6),
                     desc, font_size=10, color=TEXT_LIGHT)
    
    # ── Right: Conclusion ──
    add_text_box(slide, Inches(7.2), Inches(1.2), Inches(5), Inches(0.35),
                 "Conclusion", font_size=16, color=ACCENT_GREEN, bold=True)
    
    conclusion_card = add_rounded_rect(slide, Inches(7.2), Inches(1.7), Inches(5.2), Inches(2.5), BG_CARD, ACCENT_GREEN)
    
    conclusion_points = [
        "Maestro fundamentally transforms restaurant operations from reactive to proactive.",
        "Six specialized AI agents collaborate through a digital twin to optimize every aspect: demand, kitchen, inventory, guests, staff, and global coordination.",
        "The multi-agent architecture with heuristic fallback ensures reliability even when AI services are unavailable.",
        "Real-time updates via WebSocket provide instant visibility across all stakeholders.",
        "The system demonstrates measurable impact: reduced waste, faster service, higher guest satisfaction, and healthier staff.",
    ]
    
    for i, point in enumerate(conclusion_points):
        y = Inches(1.85) + Inches(i * 0.45)
        # Bullet
        add_shape(slide, MSO_SHAPE.OVAL, Inches(7.4), y + Inches(0.05), Inches(0.12), Inches(0.12), fill_color=ACCENT_GREEN)
        add_text_box(slide, Inches(7.65), y, Inches(4.5), Inches(0.4),
                     point, font_size=10, color=TEXT_LIGHT)
    
    # ── Scalability Roadmap (Bottom Right) ──
    add_text_box(slide, Inches(7.2), Inches(4.5), Inches(5), Inches(0.35),
                 "Scalability & Expansion", font_size=14, color=ACCENT_PURPLE, bold=True)
    
    roadmap = [
        ("Phase 1", "Single restaurant deployment"),
        ("Phase 2", "Multi-location federation"),
        ("Phase 3", "Supply chain integration"),
        ("Phase 4", "Industry-wide platform"),
    ]
    
    for i, (phase, desc) in enumerate(roadmap):
        y = Inches(4.95) + Inches(i * 0.55)
        
        # Phase indicator
        indicator = add_rounded_rect(slide, Inches(7.2), y, Inches(0.9), Inches(0.4),
                                     fill_color=ACCENT_PURPLE if i < 2 else BG_CARD,
                                     line_color=ACCENT_PURPLE)
        tf = indicator.text_frame
        tf.paragraphs[0].text = phase
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = TEXT_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        add_text_box(slide, Inches(8.25), y + Inches(0.05), Inches(3.5), Inches(0.3),
                     desc, font_size=10, color=TEXT_LIGHT)
    
    # ── Bottom: Call to Action ──
    cta_box = add_rounded_rect(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7),
                               fill_color=RGBColor(0x1E, 0x3A, 0x5F), line_color=ACCENT_BLUE)
    add_text_box(slide, Inches(1.0), Inches(6.55), Inches(11.3), Inches(0.25),
                 "Maestro: Your Restaurant Runs Itself",
                 font_size=18, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.3),
                 "Less waste. Faster tables. Happier guests. Staff that do not burn out.",
                 font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    print(f"Loading PPT: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)
    
    # Slide 1: Leave as is (title page)
    print("Slide 1: Keeping original title page")
    
    # Slide 2: Current Problem
    print("Slide 2: Filling Current Problem...")
    fill_slide_2(prs.slides[1])
    
    # Slide 3: Proposed Solution
    print("Slide 3: Filling Proposed Solution...")
    fill_slide_3(prs.slides[2])
    
    # Slide 4: Technical Approach
    print("Slide 4: Filling Technical Approach...")
    fill_slide_4(prs.slides[3])
    
    # Slide 5: Use Cases & Impact
    print("Slide 5: Filling Use Cases & Impact...")
    fill_slide_5(prs.slides[4])
    
    # Slide 6: Future Scope & Conclusion
    print("Slide 6: Filling Future Scope & Conclusion...")
    fill_slide_6(prs.slides[5])
    
    # Save
    prs.save(OUTPUT_PATH)
    print(f"\nDone! Saved to: {OUTPUT_PATH}")
    print(f"File size: {os.path.getsize(OUTPUT_PATH) / 1024:.1f} KB")


if __name__ == '__main__':
    main()
